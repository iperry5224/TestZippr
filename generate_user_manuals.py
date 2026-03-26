#!/usr/bin/env python3
"""
Generate SAELAR and SOPRA User Manuals (ISSO-focused)
Output: PowerPoint (.pptx) files compatible with Google Slides
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "user_manuals"
SCREENSHOTS = OUTPUT / "screenshots"


def add_logo_slide(prs, logo_path, title, subtitle):
    """Add title slide with logo."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1a, 0x52, 0x9c)
    bar.line.fill.background()
    
    # Logo (centered, top)
    if logo_path.exists():
        pic = slide.shapes.add_picture(str(logo_path), Inches(2.5), Inches(1.4), width=Inches(3))
        pic.left = int((prs.slide_width - pic.width) / 2)
    
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x64, 0x64, 0x64)
    p2.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, screenshot_path=None, accent_color=None):
    """Add content slide with bullet points. If screenshot_path exists, show two-column layout."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Accent bar at top
    bar_color = accent_color or RGBColor(0x1a, 0x52, 0x9c)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    
    has_screenshot = screenshot_path and Path(screenshot_path).exists()
    
    if has_screenshot:
        # Two-column: text left (4"), screenshot right (5.5")
        text_width = Inches(4)
        img_left = Inches(4.5)
        img_width = Inches(5)
        img_height = Inches(5)
    else:
        text_width = Inches(9)
        img_left = img_width = img_height = None
    
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), text_width, Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    
    # Bullets
    body_height = Inches(5.8) if has_screenshot else Inches(5.5)
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), text_width, body_height)
    tf = body.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13) if has_screenshot else Pt(14)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(4)
        p.level = 0
    
    # Screenshot (right side)
    if has_screenshot:
        slide.shapes.add_picture(str(screenshot_path), img_left, Inches(1.1), width=img_width, height=img_height)
    
    return slide


# SAELAR screenshot filenames (place in user_manuals/screenshots/saelar/)
SAELAR_SCREENSHOTS = [
    None,  # slide 0 = title
    "01_what_is.png",      # Splash or main dashboard
    "02_getting_started.png",  # Splash + sidebar with System info
    "03_main_tabs.png",    # Tab bar: NIST Assessment, Chad, Risk Calculator, etc.
    "04_nist_assessment.png",  # NIST Assessment tab, control families, Run Assessment
    "05_chad.png",         # Chad AI chat interface
    "06_ssp_poam.png",     # SSP Generator tab, POA&Ms
    "07_bod_kev.png",      # BOD 22-01 tab, Check CVEs or Dashboard
    "08_tips.png",         # Optional: sidebar with system info
    "09_need_help.png",    # Optional: AWS config or support
]


def create_saelar_manual():
    """Create SAELAR User Manual for ISSO personnel."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    logo = ASSETS / "saelar_logo.png"
    saelar_imgs = SCREENSHOTS / "saelar"
    
    add_logo_slide(prs, logo, "SAELAR User Manual",
                   "Security Architecture Evaluation & NIST 800-53 Risk Assessment – Quick Reference for ISSOs")
    
    add_content_slide(prs, "What is SAELAR?",
        [
            "SAELAR (Security Architecture Evaluation & Risk Assessment) is a cloud-based security assessment platform that helps ISSOs evaluate AWS environments against NIST 800-53 Rev 5.",
            "It uses AI (AWS Bedrock) to provide compliance guidance and generate findings.",
            "Data stays within your AWS account – no external transmission.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[1] if SAELAR_SCREENSHOTS[1] else None)
    
    add_content_slide(prs, "Getting Started",
        [
            "1. Accept the disclaimer on the splash screen.",
            "2. Configure AWS credentials (first time only): Account ID, IAM user, Access Key, Secret Key.",
            "3. Use the sidebar to enter System information: name, acronym, owner, AO, ISSO, categorization (FIPS 199).",
            "4. Select control families to assess (e.g., AC, AU, SC) and click Run Assessment.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[2] if SAELAR_SCREENSHOTS[2] else None)
    
    add_content_slide(prs, "Main Tabs – Overview",
        [
            "NIST Assessment – Run and view compliance assessments against NIST 800-53 Rev 5 controls.",
            "AWS Console – Quick link to your AWS Console for reference.",
            "Chad (AI Analyst) – Chat with AI for compliance guidance, POA&M suggestions, and risk analysis.",
            "Risk Calculator – Risk matrix, threat sources, ALE, and NIST 800-30 enhanced metrics.",
            "SSP Generator – Generate System Security Plans and POA&Ms from assessment data.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[3] if SAELAR_SCREENSHOTS[3] else None)
    
    add_content_slide(prs, "NIST Assessment Workflow",
        [
            "Select control families in the sidebar (Low / Moderate / High baseline).",
            "Click Run Assessment – SAELAR evaluates your AWS config against selected controls.",
            "Review results: Pass/Fail/Not Applicable per control, with evidence and guidance.",
            "Chad can help interpret findings and suggest remediation.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[4] if SAELAR_SCREENSHOTS[4] else None)
    
    add_content_slide(prs, "Chad AI Assistant",
        [
            "Ask Chad questions about your assessment, controls, or compliance.",
            "Request POA&M tables, risk summaries, or explanations of specific controls.",
            "Provide context by pasting control IDs or findings – Chad responds with actionable guidance.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[5] if SAELAR_SCREENSHOTS[5] else None)
    
    add_content_slide(prs, "SSP Generator & POA&Ms",
        [
            "Generate an SSP document (Word) from your assessment results.",
            "POA&Ms tab shows Plan of Action & Milestones – open findings with milestones and due dates.",
            "Use Chad to refine POA&M language before exporting.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[6] if SAELAR_SCREENSHOTS[6] else None)
    
    add_content_slide(prs, "BOD 22-01 (KEV)",
        [
            "Check CVEs against CISA Known Exploited Vulnerabilities catalog.",
            "Upload scan results or paste CVE IDs – SAELAR identifies KEVs.",
            "Generate BOD 22-01 compliance reports.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[7] if SAELAR_SCREENSHOTS[7] else None)
    
    add_content_slide(prs, "Tips for ISSOs",
        [
            "Enter system info in the sidebar before running assessments – it populates SSP and reports.",
            "Start with a few control families, then expand once familiar.",
            "Use Chad to clarify control requirements or draft remediation language.",
            "Export SSP and POA&Ms for ATO packages and continuous monitoring.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[8] if SAELAR_SCREENSHOTS[8] else None)
    
    add_content_slide(prs, "Need Help?",
        [
            "Contact your system administrator for access and configuration.",
            "AWS credentials must have permissions for Security Hub, Config, and Bedrock (see IAM requirements).",
            "Air-gapped mode: SAELAR supports local Ollama for environments without internet.",
        ], saelar_imgs / SAELAR_SCREENSHOTS[9] if SAELAR_SCREENSHOTS[9] else None)
    
    OUTPUT.mkdir(exist_ok=True)
    out_path = OUTPUT / "SAELAR_User_Manual.pptx"
    prs.save(out_path)
    return out_path


# SOPRA graphic filenames
SOPRA_GRAPHICS = [
    None, "01_what_is.png", "02_getting_started.png", "03_main_sections.png",
    "04_assessment.png", "05_dashboard.png", "06_isso_toolkit.png",
    "07_vuln_mgmt.png", "08_ssp_reports.png", "09_tips.png", "10_need_help.png",
]
SOPRA_ACCENT = RGBColor(0x00, 0xd9, 0xff)  # SOPRA cyan


def create_sopra_manual():
    """Create SOPRA User Manual for ISSO personnel."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    logo = ASSETS / "SOPRA_logo_dark.png"
    if not logo.exists():
        logo = ASSETS / "OPRA_logo_dark.png"
    
    # SOPRA title: cyan accent bar
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SOPRA_ACCENT
    bar.line.fill.background()
    if logo.exists():
        pic = slide.shapes.add_picture(str(logo), Inches(2.5), Inches(1.4), width=Inches(3))
        pic.left = int((prs.slide_width - pic.width) / 2)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "SOPRA User Manual"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    p.alignment = PP_ALIGN.CENTER
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "SAE On-Premise Risk Assessment – Quick Reference for ISSOs"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x64, 0x64, 0x64)
    p2.alignment = PP_ALIGN.CENTER
    
    sopra_imgs = SCREENSHOTS / "sopra"
    
    add_content_slide(prs, "What is SOPRA?",
        [
            "SOPRA (SAE On-Premise Risk Assessment) is an on-premise security assessment tool for enterprise infrastructure.",
            "It focuses on operational controls across 20 categories (e.g., Identity, Network, Endpoint, Physical Security).",
            "Uses CSV-based data upload – no cloud connection required for assessments.",
        ], sopra_imgs / SOPRA_GRAPHICS[1] if SOPRA_GRAPHICS[1] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "Getting Started",
        [
            "1. Open SOPRA and use the sidebar to navigate.",
            "2. Go to Assessment → Download Templates to get CSV templates for each control category.",
            "3. Populate the CSVs with your environment data (or use Demo Data for exploration).",
            "4. Upload completed CSVs in the Assessment tab to run the assessment.",
        ], sopra_imgs / SOPRA_GRAPHICS[2] if SOPRA_GRAPHICS[2] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "Main Sections",
        [
            "Dashboard – Overview of assessment results, risk score, and category status.",
            "Assessment – Upload CSVs, run assessments, view control-level results.",
            "Reports – Export assessment reports and findings.",
            "Vulnerability Management – AI-assisted remediation for identified gaps.",
            "SSP Generator – Generate System Security Plan content from assessment data.",
            "AI Assistant – Chat for guidance on controls and remediation.",
        ], sopra_imgs / SOPRA_GRAPHICS[3] if SOPRA_GRAPHICS[3] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "Assessment Workflow",
        [
            "Download Templates (Assessment tab) – Get CSV templates with required columns for each category.",
            "Fill in your data – Status, evidence, notes per control.",
            "Upload CSVs – One or more files; SOPRA aggregates and scores.",
            "Review Dashboard – Risk score, category breakdown, and findings.",
        ], sopra_imgs / SOPRA_GRAPHICS[4] if SOPRA_GRAPHICS[4] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "Dashboard Metrics",
        [
            "Categories Assessed – Number of control families evaluated.",
            "Controls Assessed – Total controls with FIPS 199 level if applicable.",
            "Risk Score – Aggregate risk percentage based on findings.",
            "Status – Pending or Complete.",
            "Click metric tiles to drill into details.",
        ], sopra_imgs / SOPRA_GRAPHICS[5] if SOPRA_GRAPHICS[5] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "ISSO Toolkit",
        [
            "Access via the ISSO Toolkit link in the sidebar.",
            "Central hub for POA&M, Risk Acceptance, STIG Import, Evidence, Incidents, Approvals.",
            "Supports continuous monitoring workflows and compliance documentation.",
        ], sopra_imgs / SOPRA_GRAPHICS[6] if SOPRA_GRAPHICS[6] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "Vulnerability Management & Remediation",
        [
            "Vulnerability Management tab – AI-assisted remediation dashboard.",
            "Review findings, get suggested fixes, track remediation status.",
            "Integration with CISA KEV and security best practices.",
        ], sopra_imgs / SOPRA_GRAPHICS[7] if SOPRA_GRAPHICS[7] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "SSP Generator & Reports",
        [
            "Generate SSP sections from your assessment results.",
            "Reports tab – Export findings in formats suitable for ATO packages.",
            "Link evidence and control status to documentation.",
        ], sopra_imgs / SOPRA_GRAPHICS[8] if SOPRA_GRAPHICS[8] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "Tips for ISSOs",
        [
            "Use Demo Data first to explore the tool without real data.",
            "Start with a few categories (e.g., Identity, Network) before full assessment.",
            "ISSO Toolkit centralizes POA&M, risk acceptance, and evidence – use it for ConMon.",
            "AI Assistant helps interpret controls and draft remediation language.",
        ], sopra_imgs / SOPRA_GRAPHICS[9] if SOPRA_GRAPHICS[9] else None, SOPRA_ACCENT)
    
    add_content_slide(prs, "Need Help?",
        [
            "Contact your system administrator for access and deployment.",
            "SOPRA can run locally or on a shared server – check your organization’s setup.",
        ], sopra_imgs / SOPRA_GRAPHICS[10] if SOPRA_GRAPHICS[10] else None, SOPRA_ACCENT)
    
    OUTPUT.mkdir(exist_ok=True)
    out_path = OUTPUT / "SOPRA_User_Manual.pptx"
    prs.save(out_path)
    return out_path


def main():
    print("Generating graphics (if needed)...")
    try:
        import create_manual_graphics
        create_manual_graphics.main()
    except Exception as e:
        print(f"Graphics: {e}")
    print("Generating user manuals...")
    saelar_path = create_saelar_manual()
    sopra_path = create_sopra_manual()
    print(f"Created: {saelar_path}")
    print(f"Created: {sopra_path}")
    print("\nTo use in Google Slides: File > Import slides > Upload > select .pptx file")


if __name__ == "__main__":
    main()
